from pptx import Presentation
import pandas as pd
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor

def addSlideTitle(prs, slideTitle):
    """ 
        Adding new slide to presentation prs with title strTitle 
		Parameters
    	----------        
        prs        : object Presentation
                     presentation where adding new slide
        slideTitle : string 
                     value of title of slide
        Returns
        ----------        
        slide      : pptx.slide.Slide
                     slide object from the python-pptx library with added title        
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    setSlideTitle(title, slideTitle)
    return slide

def setSlideTitle(title, text):
    """ 
       Set style for slide's title with value text
	   Parameters
       ----------               
       title : slide.shapes.title
               object keep info about slide's title (style, text...)
       text  : string 
               value of slide's title
       Returns
       ----------        
       title : slide.shapes.title
               title object with added formated text
    """
    title.text = text
    title.top = Inches(0.1)
    title.height = Inches(0.7)
    title.width = Inches(16)
    title.vertical_anchor = MSO_ANCHOR.TOP
    title.text_frame.paragraphs[0].font.size = Pt(32)
    return title

def addTextToSlide(slide, text, left, top, width, height, fontSize=18):
    """
       Adding text to slide
	   Parameters
       ----------               
       slide   : pptx.slide.Slide
                 slide object from the python-pptx library
       text    : string 
                 value of adding text
       left    : integer
                 Position of the left-side of the table, either as an integer in cm, or
                 as an instance of a pptx.util Length class (pptx.util.Inches for
                 example). Defaults to 4cm.
       top     : integer
                 Position of the top of the table, takes parameters as above.
       width   : integer
                 Position of the top of the table, takes parameters as above.
       height  : integer
                 Position of the top of the table, takes parameters as above.
       fontSize: integer, optional
                 font size for adding text
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    subtitle = txBox.text_frame
    subtitle.text = text
    for paragraph in subtitle.paragraphs:
        paragraph.font.size = Pt(fontSize)

def addListToSlide(slide, titleList, infoList, left, top, width, height, fontSize=18, sep='\n'):
    """
       Adding text or numeric list to slide
       Parameters
       ----------
       slide    : pptx.slide.Slide
                  slide object from the python-pptx library
       titleList: string
                  value of adding table (will display on the slide up of table)
       infoList: list
                  value of adding list (should be text of numeric values in the list)
       left     : integer
                  Position of the left-side of the table, either as an integer in cm, or
                  as an instance of a pptx.util Length class (pptx.util.Inches for
                  example).
       top      : integer
                  Position of the top of the table, takes parameters as above.
       width    : integer
                  Position of the top of the table, takes parameters as above.
       height   : integer
                  Position of the top of the table, takes parameters as above.
       fontSize : integer, optional
                  font size for adding text
       sep      : string, optional
                  contain symbol to separate value in the list (e.g. '\n' to start each member
                  of list from new line)
    """
    if titleList == '':
        text = sep.join(map(str, infoList))
    else:
        text = titleList + '\n' + sep.join(map(str, infoList))
    addTextToSlide(slide, text, left, top, width, height, fontSize)

def addTableToSlide(slide, df, left, top, width, height, fontSize=18, colWidths=[]):
    """ 
       Convert pandas dataframe to pptx table and insert it to slide with specific formats.
       Current style: white background with grey borders and black font color
	   Parameters
       ----------               
       slide    : pptx.slide.Slide
                  slide object from the python-pptx library
       df:        pd.DataFrame 
                  dataframe contain table for adding to the slide
       left     : integer
                  Position of the left-side of the table, either as an integer in cm, or
                  as an instance of a pptx.util Length class (pptx.util.Inches for
                  example). Defaults to 4cm.
       top      : integer
                  Position of the top of the table, takes parameters as above.
       width    : integer
                  Position of the top of the table, takes parameters as above.
       height   : integer
                  Position of the top of the table, takes parameters as above.
                  Use the smallest value and PPTX adjust table height by row count.
       fontSize : integer, optional
                  font size for adding text
       colWidths: list of integer, optional
                  Contain list of width of each column of the table (currently in inches)
       Returns
       ----------        
       title : shape.table
               table python pptx to manage it ourside of function               
    """                  
    # get dataframe size
    rCount = df.shape[0]
    cCount = df.shape[1]

    # set size of pptx table and location
    # use Inches(1) for make table height as small as possible
    shape = slide.shapes.add_table(rCount + 1, cCount, left, top, width, height)
    table = shape.table

    # fill first row - column names
    for c in range(cCount):
        table.cell(0, c).text = df.columns[c]

    # fill values from df
    dfValues = df.values
    for r in range(rCount):
        for c in range(cCount):
            # table.cell(r + 1, c).text = str(df.loc[r, df.columns[c]])
            table.cell(r + 1, c).text = str(dfValues[r, c])

    # set standard white format
    for r in range(rCount + 1):
        for c in range(cCount):
            cell = table.cell(r, c)
            # set border of the cell
            cell = _set_cell_border(cell)
            # set fill type to solid color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            # set color of font
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].font.size = Pt(fontSize)
            cell._tc.get_or_add_tcPr()

    # set specific column widths
    if len(colWidths) > 0:
        for col in range(len(df.columns)):
            table.columns[col].width = Inches(colWidths[col])

    return table

def SubElement(parent, tagname, **kwargs):
    """
       function for _set_cell_border
    """
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="a9a9a9", border_width='12700'):
    """
    Hack function to enable the setting of border width and border color
            (c) Steve Canny
    Parameters
      ----------        
        cell        : table.cell
                      object of one cell of the pptx table
        border_color: string, optional
                      value of RGB color of border of pptx table
        border_width: string, optional
                      value of border width of the table
    Returns
      ----------        
        cell        : table.cell
                      object of one cell of the pptx table
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
    return cell    

def addSlideWithList(prs, slideTitle, titleList, infoList, fontSize=18, sep='\n'):
    """ 
       Adding slide with title and list of values
		Parameters
    	----------        
        prs       : object Presentation
                    presentation where adding new slide
        slideTitle: string 
                    value of title of slide
        titleList : string 
                    value of adding table (will display on the slide up of table)
        infoList  : list 
                    value of adding list (should be text of numeric values in the list)
        fontSize  : integer, optional
                    font size for adding text
        sep       : string, optional
                    contain symbol to separate value in the list (e.g. '\n' to start each member
                    of list from new line)
    """   
    slide = addSlideTitle(prs, slideTitle)
    addListToSlide(slide, titleList, infoList, Inches(4), Inches(1), Inches(14), Inches(8), fontSize, sep)

def addSlideWithTable(prs, slideTitle, dfList, fontSize=12, colWidthsList=[]):
    """ 
       Adding slide with title and pandas dataframe (or dataframes)
		Parameters
    	----------        
        prs          : object Presentation
                       presentation where adding new slide
        slideTitle   : string 
                       value of title of slide
        dfList       : list
                       list of pandas dataframes  which need insert to slide
        fontSize     : integer, optional
                       font size for adding text
       colWidths: list of integer, optional
                  Contain list of width of each column of the table (currently in inches)
    """   
    slide = addSlideTitle(prs, slideTitle)
    top = 1
    for curDf in range(len(dfList)):
        df = dfList[curDf]
        # set list of column width for specific column widths
        try:
            if len(colWidthsList) > 0:
                colWidths = colWidthsList[curDf]
            else:
                colWidths = []
        except:
            colWidths = []

        if len(df) > 0:
            table = addTableToSlide(slide, df, Inches(0), Inches(top), Inches(16), Inches(1), fontSize, colWidths)
            top += 4